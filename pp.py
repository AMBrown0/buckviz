# -*- coding: utf-8 -*-
"""
Created on Fri May 29 16:05:48 2020

@author: andy
"""


import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
import seaborn as sns
from numpy import median, mean
import json, urllib.request
import sys
from pathlib import Path

from pptx import Presentation
dataFolder=Path(r"C:\\Users\\Andy.JIVEDIVE\\OneDrive - University of Buckingham\\Notes\\2020-05-22")

filename="test.pptx"
dataFile= dataFolder / filename



prs = Presentation(dataFile)
text_runs = []

for slide in prs.slides:
    print("Slide=%s" %(slide))
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
                